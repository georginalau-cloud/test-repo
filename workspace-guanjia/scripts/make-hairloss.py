#!/usr/bin/env python3
"""Build the hair-loss causation slide (PPTX) using EMU throughout."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────
GREEN_DARK  = RGBColor(0x2e, 0x7d, 0x32)
GREEN_MID   = RGBColor(0x56, 0xa3, 0x5a)
GREEN_LITE  = RGBColor(0xe8, 0xf5, 0xe9)
BLUE_DARK   = RGBColor(0x1a, 0x23, 0x7e)
BLUE_MID    = RGBColor(0x31, 0x57, 0x9e)
BLUE_LITE   = RGBColor(0xe3, 0xf0, 0xfd)
ORANGE      = RGBColor(0xe6, 0x51, 0x00)
RED_LITE    = RGBColor(0xff, 0xf3, 0xe0)
DIVIDER_CLR = RGBColor(0x88, 0x88, 0x88)
WHITE       = RGBColor(0xff, 0xff, 0xff)
BLACK       = RGBColor(0x1a, 0x1a, 0x1a)
GRAY_LITE   = RGBColor(0xf5, 0xf5, 0xf5)
GRAY_TXT    = RGBColor(0x44, 0x44, 0x44)
GRAY_MED    = RGBColor(0xdd, 0xdd, 0xdd)
ORANGE_BG   = RGBColor(0xff, 0x70, 0x43)
PATH_DOWN_BG = RGBColor(0xff, 0xeb, 0xee)
PATH_DOWN_BC = RGBColor(0xef, 0x9a, 0x9a)
PATH_DOWN_TX = RGBColor(0xc6, 0x28, 0x28)
PATH_UP_BG   = RGBColor(0xfc, 0xe4, 0xec)
PATH_UP_BC   = RGBColor(0xf4, 0x8f, 0xb1)
PATH_UP_TX   = RGBColor(0xad, 0x14, 0x57)

# ── Slide dimensions ─────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

W  = prs.slide_width   # EMU
H  = prs.slide_height # EMU
HY = H // 2           # EMU – half height

# ── Helpers ──────────────────────────────────────────────────────────────────
def rect(l, t, w, h, fill, border=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
    else:
        s.line.fill.background()
    return s

def txt(text, l, t, w, h, size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb

def pill(text, l, t, color, size=10):
    rect(l, t, Inches(1.2), Inches(0.35), color)
    txt(text, l + Inches(0.05), t + Inches(0.04), Inches(1.1), Inches(0.27),
        size=size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Helper: thin divider row ────────────────────────────────────────────────
def hdiv(y, color, l=Inches(0.5), rEdge=Inches(12.3)):
    rect(l, y, rEdge - l, Inches(0.018), color)

# ── Helper: section label ─────────────────────────────────────────────────────
def sec_lbl(text, y, color, l=Inches(0.5), rEdge=Inches(12.3)):
    txt(text, l, y, rEdge - l, Inches(0.25),
        size=9, bold=True, color=color, align=PP_ALIGN.CENTER)

# ── Helper: three-column card row ────────────────────────────────────────────
def card_row(y, h, items, border_color):
    x = Inches(0.5)
    gap = Inches(0.22)
    w  = Inches(3.9)
    for title, body in items:
        rect(x, y, w, h, WHITE, border_color)
        txt(title, x + Inches(0.12), y + Inches(0.1),  w - Inches(0.24), Inches(0.3),
            size=12, bold=True, color=GREEN_DARK)
        txt(body,  x + Inches(0.12), y + Inches(0.4), w - Inches(0.24), Inches(0.55),
            size=11, color=GRAY_TXT)
        x = x + w + gap

# ─────────────────────────────────────────────────────────────────────────────
# TOP HALF — 中医
# ─────────────────────────────────────────────────────────────────────────────
rect(Inches(0), Inches(0), W, HY, RGBColor(0xf0, 0xf7, 0xef))

pill("中医", Inches(6.1), Inches(0.12), GREEN_MID)
txt("TCM PERSPECTIVE", Inches(0), Inches(0.3), W, Inches(0.3),
    size=8, color=RGBColor(0xaa, 0xaa, 0xaa), align=PP_ALIGN.CENTER)

# Row 1
r1y = Inches(0.62); r1h = Inches(1.0)
card_row(r1y, r1h, [
    ("经络关系", "头皮归属经络\n气血不畅 → 毛孔失养"),
    ("脏腑关系", "肾虚为根本\n肝郁、湿热为标"),
    ("体质关系", "痰湿质↑ 阴虚质↑\n平和质↓（保护）"),
], GREEN_MID)
hdiv(r1y + r1h + Inches(0.08), GREEN_MID)

# Row 2: 核心病机
r2y = r1y + r1h + Inches(0.16)
sec_lbl("核心病机", r2y, GREEN_DARK)

bar_y = r2y + Inches(0.28); bar_h = Inches(0.52)
bar_defs = [
    (Inches(0.5),  Inches(2.3),  GREEN_DARK, "肾虚为本"),
    (Inches(3.2),  Inches(1.45), ORANGE,     "湿热"),
    (Inches(4.95), Inches(1.45), ORANGE,     "血瘀"),
    (Inches(6.7),  Inches(1.45), ORANGE,     "肝郁"),
    (Inches(8.9),  Inches(2.3),  GREEN_DARK, "脂溢性脱发"),
]
for l, w, fill, label in bar_defs:
    rect(l, bar_y, w, bar_h, fill)
    txt(label, l, bar_y + Inches(0.1), w, Inches(0.32),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Arrows between bars
arrow_x = [Inches(2.8), Inches(3.95), Inches(5.7), Inches(7.8)]
for ax in arrow_x:
    txt("→", ax, bar_y + Inches(0.08), Inches(0.4), bar_h,
        size=18, bold=True, color=GREEN_MID, align=PP_ALIGN.CENTER)

hdiv(bar_y + bar_h + Inches(0.14), GREEN_MID)

# Row 3: 危险因素
r3y = bar_y + bar_h + Inches(0.22)
sec_lbl("临床危险因素", r3y, GREEN_DARK)
ry2 = r3y + Inches(0.26); rh = Inches(0.72)
risk = [
    ("痰湿质",       "油腻、甜食、熬夜\n增大发病风险",    "高风险体质"),
    ("阴虚质",       "津液不足\n头皮失养",                "高风险体质"),
    ("家族遗传史",   "有家族史者\n发病风险增加",          "独立危险因素"),
]
for i, (title, body, tag) in enumerate(risk):
    x = Inches(0.5) + i * Inches(4.25)
    rw = Inches(4.0)
    rect(x, ry2, rw, rh, RED_LITE, RGBColor(0xff, 0xcc, 0x80))
    txt(title, x + Inches(0.12), ry2 + Inches(0.08), rw - Inches(0.24), Inches(0.28),
        size=13, bold=True, color=ORANGE)
    txt(body,  x + Inches(0.12), ry2 + Inches(0.34), rw - Inches(0.24), Inches(0.4),
        size=11, color=GRAY_TXT)

txt("来源：张月月 2019，北中医《脂溢性脱发的危险因素及其中医体质关系的临床研究》",
    Inches(0.3), HY - Inches(0.3), W - Inches(0.6), Inches(0.25),
    size=8, color=RGBColor(0xbb, 0xbb, 0xbb))

# ── Divider ───────────────────────────────────────────────────────────────────────
rect(Inches(0), HY, W, Inches(0.04), DIVIDER_CLR)

# ─────────────────────────────────────────────────────────────────────────────
# BOTTOM HALF — 现代医学
# ─────────────────────────────────────────────────────────────────────────────
rect(Inches(0), HY + Inches(0.04), W, HY - Inches(0.04),
     RGBColor(0xf0, 0xf4, 0xff))

pill("现代医学", Inches(5.6), HY + Inches(0.16), BLUE_MID)
txt("WESTERN MEDICINE PERSPECTIVE", Inches(0), HY + Inches(0.34), W, Inches(0.3),
    size=8, color=RGBColor(0xaa, 0xaa, 0xaa), align=PP_ALIGN.CENTER)

# Row 1: 始动因素
wm1_y = HY + Inches(0.64); wm1_h = Inches(0.88)
trigger = [
    ("二氢睾酮 DHT",    "5α-还原酶将睾酮→DHT\nDHT+雄激素受体(AR)\n启动级联→毛囊萎缩"),
    ("遗传易感基因",    "X染色体(AR,EDA2R)\n20号(PAX1/FOXA2)\n7号(HDAC9)"),
]
for i, (title, body) in enumerate(trigger):
    x = Inches(0.5) + i * Inches(6.15)
    rw = Inches(5.8)
    rect(x, wm1_y, rw, wm1_h, WHITE, BLUE_MID)
    txt(title, x + Inches(0.12), wm1_y + Inches(0.08), rw - Inches(0.24), Inches(0.28),
        size=13, bold=True, color=BLUE_DARK)
    txt(body,  x + Inches(0.12), wm1_y + Inches(0.38), rw - Inches(0.24), Inches(0.48),
        size=11, color=GRAY_TXT)

# Core path
core_y = wm1_y + wm1_h + Inches(0.14)
sec_lbl("核心病理", core_y, BLUE_DARK)
cb_y = core_y + Inches(0.25)
rect(Inches(1.5), cb_y, Inches(10.3), Inches(0.48), BLUE_DARK)
txt("DHT + AR  →  信号通路紊乱  →  毛囊微型化  →  毳毛脱落",
    Inches(1.5), cb_y + Inches(0.08), Inches(10.3), Inches(0.32),
    size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Path 3: 信号通路
path_y = cb_y + Inches(0.56)
sec_lbl("信号通路失调", path_y, BLUE_DARK)
pw2 = path_y + Inches(0.24)
path_data = [
    ("Wnt/β-catenin ↓", "毛囊再生\n信号受阻", True),
    ("Shh/Gli ↓",       "毛囊发育\n成熟受阻", True),
    ("PI3K/Akt ↓",      "干细胞互动\n毛发再生受阻", True),
    ("TGF-β/BMP ↑",     "诱导毛囊进入\n休止期", False),
]
for i, (name, desc, is_down) in enumerate(path_data):
    x = Inches(0.5) + i * Inches(3.1)
    bg = PATH_DOWN_BG if is_down else PATH_UP_BG
    bc = PATH_DOWN_BC if is_down else PATH_UP_BC
    tc = PATH_DOWN_TX if is_down else PATH_UP_TX
    rect(x, pw2, Inches(2.9), Inches(0.75), bg, bc)
    txt(name, x, pw2 + Inches(0.06), Inches(2.9), Inches(0.3),
        size=13, bold=True, color=tc, align=PP_ALIGN.CENTER)
    txt(desc, x, pw2 + Inches(0.36), Inches(2.9), Inches(0.35),
        size=10, color=GRAY_TXT, align=PP_ALIGN.CENTER)

# Row 4: 生活习惯 & 微环境
lif_y = pw2 + Inches(0.84)
sec_lbl("生活习惯 & 微环境", lif_y, BLUE_DARK)
lif_y2 = lif_y + Inches(0.24)
life = [
    ("吸烟",       "尼古丁→血管收缩\n微循环减少"),
    ("高糖高脂",   "代谢综合征\n心血管风险↑"),
    ("心理压力",   "HPA轴激活\n皮质醇↑"),
    ("睡眠不足",   "激素分泌紊乱\n急性应激↑"),
]
iw = Inches(2.7)
for i, (title, body) in enumerate(life):
    x = Inches(0.5) + i * (iw + Inches(0.2))
    rect(x, lif_y2, iw, Inches(0.65), GRAY_LITE, GRAY_MED)
    txt(title, x, lif_y2 + Inches(0.05), iw, Inches(0.25),
        size=12, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txt(body,  x, lif_y2 + Inches(0.30), iw, Inches(0.32),
        size=10, color=GRAY_TXT, align=PP_ALIGN.CENTER)

inf_x = Inches(0.5) + 4 * (iw + Inches(0.2))
rect(inf_x, lif_y2, Inches(1.45 * 2.7), Inches(0.65), ORANGE_BG)
txt("慢性炎症+氧化应激", inf_x, lif_y2 + Inches(0.04), Inches(1.45*2.7), Inches(0.25),
    size=11, color=WHITE, align=PP_ALIGN.CENTER)
txt("IL-1/TNF-α+ROS\n互促循环→毛囊损伤", inf_x, lif_y2 + Inches(0.30), Inches(1.45*2.7), Inches(0.32),
    size=10, color=WHITE, align=PP_ALIGN.CENTER)

txt("来源：中华整形外科杂志 2025年3月《雄激素性脱发发病机制的研究进展》李宇飞团队",
    Inches(0.3), H - Inches(0.3), W - Inches(0.6), Inches(0.25),
    size=8, color=RGBColor(0xbb, 0xbb, 0xbb))

out = "/Users/georginalau/Desktop/脱发的成因.pptx"
prs.save(out)
print("Saved:", out)
